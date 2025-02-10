import os
import shutil
import requests
import uuid
import re
import io
import logging
from pathlib import Path
from typing import List, Dict, Union

# Add new imports for image extraction
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation

import openai

from qdrant_client import models, QdrantClient
from qdrant_client.models import PointStruct, Distance, VectorParams
from qdrant_client.http.models import Filter, FieldCondition, MatchValue
from sentence_transformers import SentenceTransformer
from transformers import pipeline

# Re-introduce moondream
import moondream as md

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch

# Evaluation metrics imports
import nltk
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
from rouge_score import rouge_scorer
# NEW evaluation metric imports
from bert_score import score as bert_score
from sentence_transformers import util
#from bleurt import score as bleurt_score  # Corrected import
import sys
sys.path.append(r"\BARTScore")
# In questgensesh.py:
from BARTScore.bart_score import BARTScorer  # Notice the .
from bleurt import score  # Corrected import:  bleurt-pytorch

# ---------------------------
# Configuration for External Tools
# ---------------------------
POPPLER_PATH = r"c:\Users\Anandaraman\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"  # Adjust to your installation path
os.environ["PATH"] += os.pathsep + POPPLER_PATH
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# ---------------------------
# QDrant and OpenAI Setup
# ---------------------------
qdrant_client = QdrantClient(
    url="https://97752b4c-af12-4c68-953c-dcdd59810e6d.us-west-2-0.aws.cloud.qdrant.io",
    api_key="vtf2RP3Po5HXFgdFHHkM8a-ZFlTrk9bMbHWkoixtKoMVU6Rd24rgBQ"
)

api_key = "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" # Replace with your actual OpenAI API Key
os.environ["OPENAI_API_KEY"] = api_key

# ---------------------------
# Global Variables and Session ID
# ---------------------------
session_id = str(uuid.uuid4())
document_metadata = {
    "book": "Breadth first search",
    "chapter": "Chapter 1",
    "topic": "Algorithms",
    "subtopic": "Graph Theory"
}

# Create directory for extracted images
EXTRACTED_IMAGES_DIR = "extracted_images"
os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)

# ---------------------------
# Main Model Initialization
# ---------------------------
model = SentenceTransformer("all-MiniLM-L6-v2")

# ---------------------------
#  BLEURT and BARTScorer Initialization (Outside Functions - IMPORTANT)
# ---------------------------

# Initialize BLEURT scorer *once* outside any function.  This is MUCH more efficient.
#bleurt_checkpoint = "bleurt-large-512"  # Or another checkpoint like "BLEURT-20"
bleurt_scorer = score.BleurtScorer()

# Initialize BARTScorer *once* outside any function.
bart_scorer = BARTScorer(device='cpu', checkpoint='facebook/bart-large-cnn')  # Or 'cuda'



# ---------------------------
# Function Definitions
# ---------------------------
def process_file_content_text(file_path):
    """
    Processes the text content of a PDF using pdf2image and pytesseract.
    (Reverted to original, without paragraph handling)
    """
    try:
        pages = convert_from_path(file_path)
        extracted_text = []
        for page_num, page in enumerate(pages):
            text = pytesseract.image_to_string(page)
            extracted_text.append(text)
            # (Optional) Print text extraction status for each page if desired.
            # print(f"Extracted text from page {page_num}")
        return "\n".join(extracted_text)
    except Exception as e:
        logging.error(f"Error processing text with OCR: {e}")
        return ""

class ScientificInsightAnalyzer:
    def __init__(self, api_key: str):
        self.logger = logging.getLogger("ScientificInsightAnalyzer")
        self.logger.setLevel(logging.INFO)
        handler = logging.StreamHandler()
        handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        self.logger.addHandler(handler)
        self.model = md.vl(api_key=api_key)

    def analyze_image_insights(self, image: Image.Image) -> str:
        """
        Use Moondream to analyze a given image and extract diagram/graph insights.
        """
        try:
            encoded_image = self.model.encode_image(image)
            response = self.model.query(
                encoded_image,
                "Describe the key technical findings in this visualization using natural language. Focus on trends, patterns, and numerical values. Provide a single paragraph summary."
            )
            description = response["answer"]
            print(f"Moondream description: {description}")
            return description
        except Exception as e:
            self.logger.error(f"Error analyzing image: {str(e)}")
            raise

    def extract_images_from_pdf(self, pdf_path: str) -> List[Dict[str, Image.Image]]:
        """
        Extracts images from a PDF. In this basic implementation, the entire page is
        treated as an image.
        """
        extracted_images = []
        try:
            pages = convert_from_path(pdf_path)
            for page_num, page in enumerate(pages):
                temp_path = os.path.join(EXTRACTED_IMAGES_DIR, f'page_{page_num}.png')
                page.save(temp_path, 'PNG')
                print(f"Saved full page image for page {page_num} to {temp_path}")
                # Here you could add more advanced processing to detect individual figures.
                img_path = os.path.join(EXTRACTED_IMAGES_DIR, f'extracted_image_{page_num}.png')
                page.save(img_path, 'PNG')
                print(f"Extracted figure (or full page) for page {page_num} stored in: {img_path}")
                extracted_images.append({
                    "image": page,
                    "path": img_path
                })
        except Exception as e:
            self.logger.error(f"PDF Processing Error: {e}")
        return extracted_images

    def process_pdf(self, pdf_path: str) -> str:
        """
        Processes a PDF to extract insights from its diagrams/graphs.
        """
        insights = []
        extracted_images = self.extract_images_from_pdf(pdf_path)
        print(f"Total extracted images: {len(extracted_images)}")
        for idx, img_data in enumerate(extracted_images):
            try:
                insight = self.analyze_image_insights(img_data["image"])
                if insight:
                    insights.append(insight)
                    print(f"Insight for image {idx}: {insight}")
            except Exception as e:
                self.logger.error(f"Error processing image: {e}")
                continue
        return ' '.join(insights)

def process_file_content_image(file_path):
    """
    Processes diagram/graph-based content from a PDF.
    Uses Moondream (via ScientificInsightAnalyzer) to extract insights.
    """
    api_key_md = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJrZXlfaWQiOiJlODYyMDEzZC00NmVkLTRiNDYtOGMxZi0xYzYwMTUzY2M0YjkiLCJpYXQiOjE3Mzc1MjYyMjd9.0agZ8vgxwgrUJ7YMrIoBqGPs_4hsuh2zhqkwckxYkIM"
    analyzer = ScientificInsightAnalyzer(api_key_md)
    print("Starting diagram/graph extraction and insight analysis using Moondream...")
    try:
        insights = analyzer.process_pdf(file_path)
        print("Diagram insight extraction completed.")
        return insights
    except Exception as e:
        logging.error(f"Error processing PDF images: {e}")
        return ""

def clean_text(text):
    """
    Pre-process text by removing extra spaces and non-ASCII characters.
    """
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\x00-\x7F]+', '', text)
    return text.strip()

def split_text(text, chunk_size=1024, tolerance=20):
    """
    Splits text into chunks, aiming for a size close to chunk_size,
    but allowing for a tolerance and preferring to split at sentence ends.
    """
    chunks = []
    start = 0
    while start < len(text):
        # Ideal end point
        ideal_end = start + chunk_size
        # Maximum end point considering tolerance
        max_end = ideal_end + tolerance

        # If we're close to the end of the text, just take what's left
        if max_end >= len(text):
            chunks.append(text[start:])
            break

        # Find the nearest full stop within the tolerance range
        end = text.rfind('.', start + chunk_size - tolerance, max_end)

        # If no full stop is found within tolerance, force split at max_end
        if end == -1:
            end = max_end

        chunks.append(text[start:end + 1])  # Include the full stop
        start = end + 1  # Move past the full stop

    return chunks
def embed_text_chunks(chunks, model):
    """
    Embeds text chunks using a Sentence Transformer model.
    """
    embeddings = model.encode(chunks).tolist()
    print("Embedded text chunks into vector representations.")
    return embeddings

def insert_into_qdrant(collection_name, embeddings, chunks, additional_metadata: dict, batch_size=100):
    """
    Inserts text chunks (with embeddings) into the Qdrant collection.
    """
    total_points = 0
    for i in range(0, len(embeddings), batch_size):
        batch_embeddings = embeddings[i:i + batch_size]
        batch_chunks = chunks[i:i + batch_size]
        points = [
            PointStruct(
                id=str(uuid.uuid4()),  # Use UUID for unique IDs
                vector=embedding,
                payload={
                    "text": chunk,
                    "session_id": session_id, # Use passed session_id!
                    **additional_metadata
                }
            )
            for embedding, chunk in zip(batch_embeddings, batch_chunks)
        ]
        qdrant_client.upsert(collection_name=collection_name, points=points)
        total_points += len(points)
        print(f"Inserted batch of {len(points)} points into Qdrant collection '{collection_name}'.")
    print(f"Total points inserted into Qdrant: {total_points}")

def fill_placeholders(template_path, output_path, placeholders):
    """
    Fills in placeholders in a template file and writes the updated content to an output file.
    """
    with open(template_path, 'r', encoding='utf-8') as file:
        template = file.read()
    for placeholder, value in placeholders.items():
        template = template.replace(f"{{{placeholder}}}", value)
    with open(output_path, 'w', encoding='utf-8') as file:
        file.write(template)
    print(f"Filled placeholders in template and saved updated file to {output_path}")

def get_gpt_response(system_prompt, user_prompt):
    """
    Gets a response from GPT (using the gpt-3.5-turbo model).
    """
    print("Sending prompt to GPT for generation...")
    openai.api_key = api_key
    client = openai.OpenAI()
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
    )
    gpt_response = response.choices[0].message.content
    print("Received response from GPT.")
    return gpt_response

def find_topics_and_generate_hypothetical_text(hypothetical_prompt_path, updated_hypothetical_path, academic_level, major, course_name, taxonomy_level, topics):
    """
    Generates hypothetical text based on topics.  (This function is now used.)
    """
    with open(hypothetical_prompt_path, "r", encoding="utf8") as file:
        hypothetical_prompt = file.read()  # Read the original prompt
    placeholders = {
        "course_name": course_name,
        "academic_level": academic_level,
        "topics": topics,
        "major": major,
        "taxonomy_level": taxonomy_level
    }
    fill_placeholders(hypothetical_prompt_path, updated_hypothetical_path, placeholders)

    with open(updated_hypothetical_path, "r", encoding="utf8") as file:
        updated_hypothetical_prompt = file.read() # Read *updated* prompt

    print("Generating hypothetical text using updated prompt...")
    return get_gpt_response("You are a helpful assistant skilled at analyzing text.", updated_hypothetical_prompt)

def view_points_in_collection(collection_name: str, limit: int = 1000):
    points = qdrant_client.scroll(
        collection_name=collection_name,
        scroll_filter=None,
        limit=limit
    )
    print(f"Retrieved {len(points[0])} points from Qdrant collection '{collection_name}'.")  # Corrected line
    return points[0]

def search_results_from_qdrant(qdrant_client, collection_name, embedded_vector, limit=25, session_id_filter=None):
    """
    Searches in Qdrant using an embedded vector.
    """
    query_filter = None
    if session_id_filter:
        query_filter = Filter(must=[FieldCondition(key="session_id", match=MatchValue(value=session_id_filter))])
    else: #Added this
        query_filter = None

    search_results = qdrant_client.search(
        collection_name=collection_name,
        query_vector=embedded_vector.tolist(),
        limit=limit,
        query_filter=query_filter
    )
    print(f"Search in collection '{collection_name}' returned {len(search_results)} results.")
    return search_results

def evaluate_generated_questions(candidate: str, reference: str):
    """
    Compute various evaluation metrics for the generated questions.
    Returns a dictionary of scores.
    """
    # Tokenize sentences for BLEU (if needed - depends on your questions)
    candidate_sentences = [s.strip() for s in candidate.split('.') if s.strip()]
    reference_sentences = [s.strip() for s in reference.split('.') if s.strip()]

    # Compute BLEU-4
    smoothie = SmoothingFunction().method4
    bleu_scores = [
        sentence_bleu([reference_sentences], cand.split(), smoothing_function=smoothie)
        for cand in candidate_sentences
    ]
    avg_bleu = sum(bleu_scores) / len(bleu_scores) if bleu_scores else 0

    # Compute ROUGE scores
    scorer = rouge_scorer.RougeScorer(['rouge1', 'rouge2', 'rougeL'], use_stemmer=True)
    rouge_scores = scorer.score(reference, candidate)

    # Compute BERTScore
    _, _, bertscore_f1 = bert_score([candidate], [reference], lang="en", rescale_with_baseline=True)
    bertscore_f1 = bertscore_f1.mean().item()

    # Compute BLEURT - Use pre-initialized scorer
    bleurt_scores_list = bleurt_scorer.score(references=[reference], candidates=[candidate])
    bleurt_score = sum(bleurt_scores_list) / len(bleurt_scores_list) if bleurt_scores_list else 0


    # Compute BARTScore - Use pre-initialized scorer
    bart_scores_list = bart_scorer.score([reference], [candidate], batch_size=4)  # Returns a list
    bartscore = sum(bart_scores_list) / len(bart_scores_list) if bart_scores_list else 0



    # Compute Q-BLEU4 (custom BLEU4 for question generation)
    q_bleu4 = bleu_scores[-1] if bleu_scores else 0

    # Compute QSTS (semantic similarity using Sentence-BERT)
    qsts_score = util.pytorch_cos_sim(model.encode(candidate), model.encode(reference)).item()


    return {
        "BLEU-4": avg_bleu,
        "ROUGE-L": rouge_scores['rougeL'].fmeasure,
        "BERTScore": bertscore_f1,
        "BLEURT": bleurt_score,
        "BARTScore": bartscore,
        "Q-BLEU4": q_bleu4,
        "QSTS": qsts_score,
    }

def regenerate_questions_if_needed(current_questions: str, reference: str, final_prompt: str, thresholds: dict):
    """
    Evaluate generated questions and, if below threshold, regenerate with feedback.
    Returns final accepted questions and feedback.
    """
    scores = evaluate_generated_questions(current_questions, reference)
    print("Evaluation Metrics:", scores)
    below_threshold = False
    feedback_comments = []

    # Check each metric against its threshold
    for metric, score in scores.items():
        if score < thresholds.get(metric, 0):
            below_threshold = True
            feedback_comments.append(f"{metric} score is {score:.2f}, which is below the threshold of {thresholds[metric]:.2f}.")

    if below_threshold:
        # Append feedback into prompt for regeneration
        feedback_text = " ".join(feedback_comments) + " Please regenerate the questions with improvements addressing these issues."
        augmented_prompt = final_prompt + "\n\n" + feedback_text
        new_questions = get_gpt_response("You are a helpful assistant skilled at automatic question generation.", augmented_prompt)
        # Optionally, one might run a loop until scores are acceptable.
        return new_questions, feedback_text, scores
    else:
        return current_questions, "All evaluation metrics are above thresholds.", scores

def generate_questions(final_user_prompt_path, updated_final_user_prompt_path, retrieved_context, course_name, num_questions, academic_level, taxonomy, topics_list, major, evaluation_prompt_path, updated_evaluation_prompt_path, thresholds):
    """Generates questions based on retrieved context and user prompt."""

    # Define Bloom's Taxonomy descriptions
    blooms_taxonomy_descriptions = """
    Please find the explanation of each level of the Bloom's taxonomy:
    Remember: retrieve, recall, or recognize relevant knowledge from long-term memory.
    Understand: demonstrate comprehension through one or more forms of explanation.
    Apply: use information or a skill in a new situation.
    Analyze: break material into its constituent parts and determine how the parts relate to one another and/or to an overall structure or purpose.
    Evaluate: make judgments based on criteria and standards.
    Create: put elements together to form a new coherent or functional whole; reorganize elements into a new pattern or structure.
    """

    # Fill placeholders in the final user prompt
    placeholders_final = {
        "content": retrieved_context,
        "num_questions": num_questions,
        "course_name": course_name,
        "taxonomy": taxonomy,
        "major": major,
        "academic_level": academic_level,
        "topics_list": topics_list,
        "blooms_taxonomy_descriptions": blooms_taxonomy_descriptions,  # Add Bloom's descriptions
    }
    fill_placeholders(final_user_prompt_path, updated_final_user_prompt_path, placeholders_final)


    with open(updated_final_user_prompt_path, "r", encoding="utf8") as file:
        final_user_prompt = file.read()

    # Generate initial questions
    print("Generating initial questions using final prompt...")
    initial_questions = get_gpt_response("You are a helpful assistant skilled at automatic question generation.", final_user_prompt)
    print("Initial questions generated.")

    # Fill in placeholders for the evaluation prompt
    placeholders_evaluation = {
        "academic_level": academic_level,
        "course_name": course_name,
        "major": major,
        "topics_list": topics_list,
        "taxonomy_level": taxonomy,
        "question_content": initial_questions
    }

    fill_placeholders(evaluation_prompt_path, updated_evaluation_prompt_path, placeholders_evaluation)

    with open(updated_evaluation_prompt_path, 'r', encoding='utf-8') as f:
        evaluation_prompt_content = f.read()

    regenerated_questions, feedback, scores = regenerate_questions_if_needed(
    initial_questions, retrieved_context, final_user_prompt, thresholds
)
    return regenerated_questions, feedback, scores  # Return questions, feedback, and scores


# ---------------------------
# Main Code Execution
# ---------------------------
def main():
    pdf_path = r"content/BFS_notespdf.pdf"

    # ----------- Text Extraction -----------
    print("Starting text extraction (text-only)...")
    text_content = process_file_content_text(pdf_path) # Reverted function
    cleaned_text = clean_text(text_content)
    print("Text extraction completed.")

     # ----------- Diagram/Graph/Illustration Extraction & Insight Analysis -----------
    print("Starting diagram extraction and insight analysis...")
    diagram_insights = process_file_content_image(pdf_path)
    cleaned_diagram_insights = clean_text(diagram_insights)
    print("Cleaned diagram insights:", cleaned_diagram_insights)

    # Combine the OCR text with diagram insights.
    combined_text = cleaned_text + " " + cleaned_diagram_insights
    print(f"Combined text length: {len(combined_text)} characters.")


    # --- Embedding & Storing in Qdrant (Text Collection Only) -----------
    print("Splitting combined text into chunks...")
    text_chunks = split_text(combined_text, chunk_size=1024, tolerance=20)  # Use split_text
    print("Embedding text chunks...")
    text_embeddings = embed_text_chunks(text_chunks, model)
    print("Inserting text chunks into Qdrant...")
    insert_into_qdrant("qgen", text_embeddings, text_chunks, document_metadata)
    print("Text and diagram insights processing and storage completed.")

    # --- Define parameters (can be from user input or config) ---
    course_name = "Data Structures and Algorithms"
    num_questions = "15"
    academic_level = "undergraduate"
    taxonomy_level = "Create"
    topics_list = "Breadth First Search, Shortest Path"  # Comma-separated
    major = "Computer Science"
    topics = topics_list  # For the hypothetical text generation

     # --- Define Thresholds ---
    thresholds = {
    "BLEU-4": 0.1,
    "ROUGE-L": 0.2,
    "BERTScore": 0.4,
    "BLEURT": 0.3,
    "BARTScore": -0.5,
    "Q-BLEU4": 0.1,
    "QSTS": 0.3,
}


    # --- Hypothetical Sub-topic Generation ---
    hypothetical_prompt_path = r"content/hypothetical_prompt.txt"
    updated_hypothetical_path = r"content/updated_hypothetical_prompt.txt"

    hypothetical_text = find_topics_and_generate_hypothetical_text(hypothetical_prompt_path, updated_hypothetical_path, academic_level, major, course_name, taxonomy_level, topics)
    print("Generated Hypothetical Text (Sub-topics):")
    print(hypothetical_text)

    # --- Search Qdrant using Hypothetical Text ---
    print(f"Performing a search in Qdrant for query: '{hypothetical_text}'")
    query_embedding = model.encode(hypothetical_text)
     # Pass session_id to filter results
    search_results = search_results_from_qdrant(qdrant_client, "qgen", query_embedding, limit=15, session_id_filter=session_id)
    print("Search Results:")
    retrieved_context = ""
    for result in search_results:
        text_snippet = result.payload.get("text", "")
        print(f"ID: {result.id}, Score: {result.score:.4f}, Snippet: {text_snippet[:200]}")  # Show first 200 chars
        retrieved_context += text_snippet + " "  # Accumulate context

    # --- Question Generation ---

    user_prompt_path = r"content/user_prompt.txt"  # Original user prompt
    updated_user_prompt_path = r"content/updated_user_prompt.txt" # Updated Original user prompt
    final_user_prompt_path = r"content/final_user_prompt.txt" # Final user prompt
    updated_final_user_prompt_path = r"content/updated_final_user_prompt.txt" #Updated final user prompt
    evaluation_prompt_path = r"content/evaluation_prompt.txt"
    updated_evaluation_prompt_path = r"content/updated_evaluation_prompt.txt"

    #First fill the original user prompt
    placeholders_user = {
        "course_name": course_name,
        "num_questions": num_questions,
        "academic_level": academic_level,
        "topics_list": topics_list,  # This is correct
        "taxonomy": taxonomy_level,
        "major": major
    }
    fill_placeholders(user_prompt_path, updated_user_prompt_path, placeholders_user)

    with open(updated_user_prompt_path, "r", encoding="utf8") as file:
        updated_user_prompt = file.read() #Storing filled user prompt to use it later for final prompt

    generated_questions, evaluation_feedback, scores  = generate_questions(
        final_user_prompt_path, updated_final_user_prompt_path, retrieved_context,
        course_name, num_questions, academic_level, taxonomy_level,
        topics_list, major, evaluation_prompt_path, updated_evaluation_prompt_path, thresholds
    )

    print("\nGenerated Questions:")
    print(generated_questions)
    print("\nEvaluation Feedback:")
    print(evaluation_feedback)
    print("\nEvaluation scores:")
    print(scores)


if __name__ == "__main__":
    main()