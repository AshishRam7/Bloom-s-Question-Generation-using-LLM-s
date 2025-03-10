# ---------------------------
# Import Statements
# ---------------------------
import os
import shutil
import requests
import json 
import uuid
import re
import io
import logging
from pathlib import Path
from typing import List, Dict, Union
from datetime import datetime
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
from langchain_community.document_loaders import PyMuPDFLoader
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
from docx import Document
from pptx import Presentation
import openai
from qdrant_client import models, QdrantClient
from qdrant_client.models import PointStruct, Distance, VectorParams
from qdrant_client.http.models import Filter, FieldCondition, MatchValue
from sentence_transformers import SentenceTransformer
from transformers import pipeline
import moondream as md
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch
import nltk
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
from rouge_score import rouge_scorer
from bert_score import score as bert_score
from sentence_transformers import util
import sys
sys.path.append(r"\BARTScore")
from BARTScore.bart_score import BARTScorer
from bleurt.bleurt import score

# ---------------------------
# Configuration for External Tools
# ---------------------------
POPPLER_PATH = r"C:\Users\VBALA\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin" # Adjust to your installation path
os.environ["PATH"] += os.pathsep + POPPLER_PATH
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ---------------------------
# QDrant and OpenAI Setup
# ---------------------------
qdrant_client = QdrantClient(
    url="https://97752b4c-af12-4c68-953c-dcdd59810e6d.us-west-2-0.aws.cloud.qdrant.io",
    api_key="vtf2RP3Po5HXFgdFHHkM8a-ZFlTrk9bMbHWkoixtKoMVU6Rd24rgBQ",
    timeout=10000
)

api_key = "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" # Replace with your actual OpenAI API Key
os.environ["OPENAI_API_KEY"] = api_key #Still setting openai key although not used directly in get_gpt_response

# ---------------------------
# Global Variables and Session ID
# ---------------------------
session_id = str(uuid.uuid4())
document_metadata = {
    "book": "Algorithms",
    "chapter": "Chapter 1",
    "topic": "Graph Algorithms",
    "subtopic": "Search Algorithms"
}

# Create directory for extracted images
EXTRACTED_IMAGES_DIR = "extracted_images"
os.makedirs(EXTRACTED_IMAGES_DIR, exist_ok=True)

# ---------------------------
# Main Model Initialization
# ---------------------------
model = SentenceTransformer("all-MiniLM-L6-v2")


# ---------------------------
# Feedback Handling System
# ---------------------------
class FeedbackHandler:
    def __init__(self):
        self.feedback_history = []

    def collect_feedback(self, automatic_feedback: str, user_feedback: str = None) -> str:
        """Combine automatic metrics with user feedback"""
        combined = f"Automatic Evaluation:\n{automatic_feedback}"
        if user_feedback:
            combined += f"\n\nUser Feedback:\n{user_feedback}"
        self.feedback_history.append(combined)
        return combined

    def generate_iteration_prompt(self, base_prompt: str, feedback: str) -> str:
        """Augment original prompt with feedback"""
        return f"""{base_prompt}

IMPROVEMENT INSTRUCTIONS:
- Carefully analyze the following feedback
- Address each point systematically
- Maintain original question requirements
- Explicitly state improvements made

FEEDBACK:
{feedback}"""


# ---------------------------
# Session Management
# ---------------------------
class SessionManager:
    def __init__(self, qdrant_client):
        self.client = qdrant_client
        self.current_session = {
            "session_id": session_id, # Use global session_id
            "timestamps": [],
            "interactions": []
        }
        # Initialize session history collection if it doesn't exist (optional, you might do this elsewhere)
        # try:
        #     self.client.get_collection(collection_name="session_history")
        # except Exception as e: # Replace with specific exception if Qdrant has one for collection not found
        #     print(f"Session history collection not found, creating...")
        #     self.client.create_collection(
        #         collection_name="session_history",
        #         vectors_config=VectorParams(size=384, distance=Distance.COSINE), # Adjust vector size if needed
        #     )
        #     print("Session history collection created.")

    def log_interaction(self, interaction_type: str, data: dict):
        """Store session context in Qdrant"""
        record = {
            "timestamp": datetime.now().isoformat(),
            "type": interaction_type,
            "data": data
        }

        self.client.upsert(
            collection_name="session_history",
            points=[PointStruct(
                id=str(uuid.uuid4()),
                vector=model.encode(json.dumps(record)).tolist(), # Embed the interaction record
                payload={
                    "session_id": self.current_session["session_id"], # Store session ID for filtering
                    "interaction": record
                }
            )]
        )
        print(f"Logged interaction: {interaction_type}")

    def get_session_context(self, lookback=3):
        """Retrieve recent session context"""
        # This part might need adjustment based on how you want to query and use context.
        # For now, it's a placeholder for context retrieval logic.
        print("Session context retrieval is not fully implemented in this example.")
        return None # Placeholder for now


# ----------------------------------------
# Extraction of Textual Data from Document
# ----------------------------------------

def extract_tables(pdf_path):
    tables = {}
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            extracted_tables = page.extract_tables()

            for table_idx, table in enumerate(extracted_tables):
                if table and len(table) > 1:
                    headers = table[0] if None not in table[0] else [f"Column_{j+1}" for j in range(len(table[1]))]
                    data = table[1:]

                    if data:
                        tables.setdefault(i + 1, []).append((headers, data))

    return tables

def extract_text(pdf_path):
    loader = PyMuPDFLoader(pdf_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def clean_text(text):
    """
    Pre-process text by removing extra spaces and non-ASCII characters.
    """
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\x00-\x7F]+', '', text)
    return text.strip()

# ----------------------------------------
# Extraction of Visual Data from Document
# ----------------------------------------

class ScientificInsightAnalyzer:
    def __init__(self, api_key: str):
        self.logger = logging.getLogger("ScientificInsightAnalyzer")
        self.logger.setLevel(logging.INFO)
        handler = logging.StreamHandler()
        handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        self.logger.addHandler(handler)
        self.model = md.vl(api_key=api_key)

    def analyze_image_insights(self, image: Image.Image) -> str:
        """
        Use Moondream to analyze a given image and extract diagram/graph insights.
        """
        try:
            encoded_image = self.model.encode_image(image)
            response = self.model.query(
                encoded_image,
                "Describe the key technical findings in this visualization using natural language. Focus on trends, patterns, and numerical values. Provide a single paragraph summary. If it doesn't hold any significance conceptual information, do not generate unnecessary insights"
            )
            description = response["answer"]
            #print(f"Moondream description: {description}")
            return description
        except Exception as e:
            self.logger.error(f"Error analyzing image: {str(e)}")
            raise

    def extract_images_from_pdf(self, pdf_path: str) -> List[Dict[str, Image.Image]]:
        doc = fitz.open(pdf_path)
        images = []

        for i, page in enumerate(doc):
            for img_index, img in enumerate(page.get_images(full=True)):  # Get all images on the page
                xref = img[0]  # Image XREF
                pix = fitz.Pixmap(doc, xref)

                # Convert CMYK and other unsupported colorspaces to RGB
                if pix.colorspace.n > 3:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                img_filename = f"image_page_{i+1}img{img_index+1}.png"
                img_path = os.path.join(EXTRACTED_IMAGES_DIR, img_filename)
                pix.save(img_path)
                images.append(img_path)
        return images

    def process_pdf(self, pdf_path: str) -> str:
        """
        Processes a PDF to extract insights from its diagrams/graphs.
        """
        insights = []
        extracted_images = self.extract_images_from_pdf(pdf_path) # Still using pdf2image based extraction for Moondream
        print(f"Total extracted images: {len(extracted_images)}")
        for idx, img_path in enumerate(extracted_images):
            try:
                if isinstance(img_path, str):  # If it's a file path
                    image = Image.open(img_path)  # Open image using PIL
                    insight = self.analyze_image_insights(image)
                else:
                    raise ValueError(f"Unexpected format for image data: {img_path}")
                
                if insight:
                    insights.append(insight)
                    print(f"Insight for image {idx}: {insight}\n")
            except Exception as e:
                self.logger.error(f"Error processing image {idx}: {e}")
                continue
        return ' '.join(insights)

def process_file_content_image(file_path):
    api_key_md = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJrZXlfaWQiOiJlODYyMDEzZC00NmVkLTRiNDYtOGMxZi0xYzYwMTUzY2M0YjkiLCJpYXQiOjE3Mzc1MjYyMjd9.0agZ8vgxwgrUJ7YMrIoBqGPs_4hsuh2zhqkwckxYkIM"
    analyzer = ScientificInsightAnalyzer(api_key_md)
    print("Starting diagram/graph extraction and insight analysis using Moondream...")
    try:
        insights = analyzer.process_pdf(file_path)
        print("Diagram insight extraction completed.")
        return insights
    except Exception as e:
        logging.error(f"Error processing PDF images: {e}")
        return ""


# ---------------------------------------------
# Chunking, Embedding and Inserting into Qdrant
# ---------------------------------------------

def split_text(text, chunk_size=1024, tolerance=20):
    chunks = []
    start = 0
    while start < len(text):
        # Ideal end point
        ideal_end = start + chunk_size
        # Maximum end point considering tolerance
        max_end = ideal_end + tolerance

        # If we're close to the end of the text, just take what's left
        if max_end >= len(text):
            chunks.append(text[start:])
            break

        # Find the nearest full stop within the tolerance range
        end = text.rfind('.', start + chunk_size - tolerance, max_end)

        # If no full stop is found within tolerance, force split at max_end
        if end == -1:
            end = max_end

        chunks.append(text[start:end + 1])  # Include the full stop
        start = end + 1  # Move past the full stop

    return chunks

def embed_text_chunks(chunks, model):
    embeddings = model.encode(chunks).tolist()
    print("Embedded text chunks into vector representations.")
    return embeddings

def insert_into_qdrant(collection_name, embeddings, chunks, additional_metadata: dict, batch_size=100):
    total_points = 0
    for i in range(0, len(embeddings), batch_size):
        batch_embeddings = embeddings[i:i + batch_size]
        batch_chunks = chunks[i:i + batch_size]
        points = [
            PointStruct(
                id=str(uuid.uuid4()),  # Use UUID for unique IDs
                vector=embedding,
                payload={
                    "text": chunk,
                    "session_id": session_id, # Use passed session_id!
                    **additional_metadata
                }
            )
            for embedding, chunk in zip(batch_embeddings, batch_chunks)
        ]
        qdrant_client.upsert(collection_name=collection_name, points=points)
        total_points += len(points)
        # print(f"Inserted batch of {len(points)} points into Qdrant collection '{collection_name}'.")
    print(f"Total points inserted into Qdrant: {total_points}")



# --------------
# Util Functions
# --------------
def fill_placeholders(template_path, output_path, placeholders):
    with open(template_path, 'r', encoding='utf-8') as file:
        template = file.read()
    for placeholder, value in placeholders.items():
        template = template.replace(f"{{{placeholder}}}", value)
    with open(output_path, 'w', encoding='utf-8') as file:
        file.write(template)
    print(f"Filled placeholders in template and saved updated file to {output_path}")


def view_points_in_collection(collection_name: str, limit: int = 1000):
    points = qdrant_client.scroll(
        collection_name=collection_name,
        scroll_filter=None,
        limit=limit
    )
    #print(f"Retrieved {len(points[0])} points from Qdrant collection '{collection_name}'.") 
    return points[0]

def search_results_from_qdrant(qdrant_client, collection_name, embedded_vector, limit=25, session_id_filter=None):
    """
    Searches in Qdrant using an embedded vector.
    """
    query_filter = None
    if session_id_filter:
        query_filter = Filter(must=[FieldCondition(key="session_id", match=MatchValue(value=session_id_filter))])
    else: #Added this
        query_filter = None

    search_results = qdrant_client.search(
        collection_name=collection_name,
        query_vector=embedded_vector.tolist(),
        limit=limit,
        query_filter=query_filter
    )
    #print(f"Search in collection '{collection_name}' returned {len(search_results)} results.")
    return search_results

def _is_improvement(new_scores, old_scores, threshold=0.05):
    """Check if at least 2 metrics improved by threshold"""
    improvements = 0
    for k in new_scores:
        if new_scores[k] > old_scores[k] + threshold: # Changed to strictly greater than threshold
            improvements += 1
    return improvements >= 2

# ---------------------------------------------
# Function to Get Response from Open-source LLMs
# ---------------------------------------------

def get_llm_response(system_prompt, user_prompt,llm="gemini"):
    if llm=="gemini":
        print("Sending prompt to Gemini for generation...")
        gemini_api_key = "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" # Replace with your actual Deepseek API key
        headers = {
            "Authorization": f"Bearer {gemini_api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "<YOUR_SITE_URL>",  # Optional
            "X-Title": "<YOUR_SITE_NAME>",  # Optional
        }
        data = json.dumps({
            "model": "google/gemini-2.0-flash-thinking-exp:free",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
        })

        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
            data=data,
        )

        if response.status_code == 200:
            gemini_response = response.json()["choices"][0]["message"]["content"]
            print("Received response from Gemini.")
            return gemini_response
        else:
            print(f"Error: Gemini API request failed with status code {response.status_code}")
            print(f"Response text: {response.text}") # Print full response for debugging
            return "Error: Failed to get response from Gemini."

    elif llm=="gpt3.5":
        print("Sending prompt to GPT for generation...\n")
        client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))  # Create a client
        response = client.chat.completions.create( 
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        gpt_response = response.choices[0].message.content
        print("Received response from GPT.\n")
        return gpt_response

    elif llm=="deepseek":
        print("Sending prompt to Deepseek for generation...")
        deepseek_api_key = "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" 
        headers = {
            "Authorization": f"Bearer {deepseek_api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "<YOUR_SITE_URL>",  # Optional
            "X-Title": "<YOUR_SITE_NAME>",  # Optional
        }
        data = json.dumps({
            "model": "deepseek/deepseek-r1:free",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
        })
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
            data=data,
        )
        if response.status_code == 200:
            deepseek_response = response.json()["choices"][0]["message"]["content"]
            print("Received response from Deepseek.")
            return deepseek_response
        else:
            print(f"Error: Deepseek API request failed with status code {response.status_code}")
            print(f"Response text: {response.text}")
            return "Error: Failed to get response from Deepseek."

    else:
        return "Invalid llm details\n"

#--------------------------------
# Generation of Hypothetical Text
#--------------------------------

def find_topics_and_generate_hypothetical_text(hypothetical_prompt_path, updated_hypothetical_path, academic_level, major, course_name, taxonomy_level, topics,llm):
    with open(hypothetical_prompt_path, "r", encoding="utf8") as file:
        hypothetical_prompt = file.read() 
    placeholders = {
        "course_name": course_name,
        "academic_level": academic_level,
        "topics": topics,
        "major": major,
        "taxonomy_level": taxonomy_level
    }
    fill_placeholders(hypothetical_prompt_path, updated_hypothetical_path, placeholders)

    with open(updated_hypothetical_path, "r", encoding="utf8") as file:
        updated_hypothetical_prompt = file.read()

    print("Generating hypothetical text using updated prompt...\n")
    return get_llm_response("You are a helpful assistant skilled at analyzing text.", updated_hypothetical_prompt,llm)

#---------------------------------------
# Generation and Evaluation of Questions
#---------------------------------------
def generate_questions(final_user_prompt_path, updated_final_user_prompt_path, retrieved_context, course_name, num_questions, academic_level, taxonomy, topics_list, major, evaluation_prompt_path, updated_evaluation_prompt_path, thresholds, feedback_handler, final_user_prompt, llm): 
    blooms_taxonomy_descriptions = """
    Please find the explanation of each level of the Bloom's taxonomy:
    Remember: retrieve, recall, or recognize relevant knowledge from long-term memory.
    Understand: demonstrate comprehension through one or more forms of explanation.
    Apply: use information or a skill in a new situation.
    Analyze: break material into its constituent parts and determine how the parts relate to one another and/or to an overall structure or purpose.
    Evaluate: make judgments based on criteria and standards.
    Create: put elements together to form a new coherent or functional whole; reorganize elements into a new pattern or structure.
    """

    placeholders_final = {
        "content": retrieved_context,
        "num_questions": num_questions,
        "course_name": course_name,
        "taxonomy": taxonomy,
        "major": major,
        "academic_level": academic_level,
        "topics_list": topics_list,
        "blooms_taxonomy_descriptions": blooms_taxonomy_descriptions,  # Add Bloom's descriptions
    }
    fill_placeholders(final_user_prompt_path, updated_final_user_prompt_path, placeholders_final)


    with open(updated_final_user_prompt_path, "r", encoding="utf8") as file:
        final_user_prompt_for_generation = file.read() # Use a separate variable to avoid confusion

    # Generate initial questions
    print("Generating initial questions using final prompt...\n")
    initial_questions = get_llm_response("You are a helpful assistant skilled at automatic question generation.", final_user_prompt_for_generation,llm) # Use the filled prompt for generation
    print("Initial questions generated:\n")
    print(initial_questions)

    # Fill in placeholders for the evaluation prompt
    placeholders_evaluation = {
        "academic_level": academic_level,
        "course_name": course_name,
        "major": major,
        "topic_list": topics_list,
        "taxonomy_level": taxonomy,
        "question_content": initial_questions
    }

    fill_placeholders(evaluation_prompt_path, updated_evaluation_prompt_path, placeholders_evaluation)

    with open(updated_evaluation_prompt_path, 'r', encoding='utf-8') as f:
        evaluation_prompt_content = f.read()

    print("\nLLM-guided Evaluation Results:\n")
    llm_eval_results=get_llm_response("You are a skilled assistant at evaluating LLM-generated questions.",evaluation_prompt_content,llm)
    print(llm_eval_results)

    
    # Use the new regeneration function with feedback handler
    regenerated_questions, feedback, scores = regenerate_questions_with_feedback(
        initial_questions, retrieved_context, final_user_prompt, feedback_handler, thresholds, llm, llm_eval_results # Pass feedback_handler and original final_user_prompt
    )
    return regenerated_questions, feedback, scores  # Return questions, feedback, and scores

def regenerate_questions_with_feedback(current_questions: str, reference: str,base_prompt: str, feedback_handler: FeedbackHandler, thresholds: dict,llm, user_feedback: str = None, max_retries=3):
    """Enhanced regeneration with combined feedback"""
    best_questions = current_questions
    best_scores = evaluate_generated_questions(current_questions, reference) # Evaluate initial questions to have a starting point
    for attempt in range(3):
        scores = evaluate_generated_questions(best_questions, reference)
        auto_feedback = "\n".join([f"{k}: {v:.2f}" for k,v in scores.items()])
        combined_feedback = feedback_handler.collect_feedback(auto_feedback, user_feedback)

        augmented_prompt = feedback_handler.generate_iteration_prompt(base_prompt, combined_feedback)
        new_questions = get_llm_response("Question Generation Expert", augmented_prompt,llm)

        new_scores = evaluate_generated_questions(new_questions, reference)
        if _is_improvement(new_scores, scores):
            best_questions = new_questions
            best_scores = new_scores
            print(f"Improvement found in attempt {attempt+1}. Keeping new questions.")
        else:
            print(f"No significant improvement in attempt {attempt+1}. Stopping regeneration.")
            break # Stop if no improvement

    return best_questions, combined_feedback, best_scores

def evaluate_generated_questions(candidate: str, reference: str):
    candidate_sentences = [s.strip() for s in candidate.split('.') if s.strip()]
    reference_sentences = [s.strip() for s in reference.split('.') if s.strip()]

    # Compute BLEU-4
    smoothie = SmoothingFunction().method4
    bleu_scores = [
        sentence_bleu([reference_sentences], cand.split(), smoothing_function=smoothie)
        for cand in candidate_sentences
    ]
    avg_bleu = sum(bleu_scores) / len(bleu_scores) if bleu_scores else 0

    # Compute QSTS (semantic similarity using Sentence-BERT)
    qsts_score = util.pytorch_cos_sim(model.encode(candidate), model.encode(reference)).item()


    return {
        "BLEU-4": avg_bleu,
        "QSTS": qsts_score,
    }

def regenerate_questions_if_needed(current_questions: str, reference: str, final_prompt: str, thresholds: dict,llm, llm_eval_results):

    scores = evaluate_generated_questions(current_questions, reference)
    print("Evaluation Metrics:", scores)
    below_threshold = False
    feedback_comments = []

    # Check each metric against its threshold
    for metric, score in scores.items():
        if score < thresholds.get(metric, 0):
            below_threshold = True
            feedback_comments.append(f"{metric} score is {score:.2f}, which is below the threshold of {thresholds[metric]:.2f}.")

    if below_threshold:
        # Append feedback into prompt for regeneration
        feedback_text = " ".join(feedback_comments) + " Please regenerate the questions with improvements addressing these issues."
        augmented_prompt = final_prompt + "\n\n" + feedback_text + "\n\n"+ llm_eval_results
        new_questions = get_llm_response("You are a helpful assistant skilled at automatic question generation.", augmented_prompt,llm)
        # Optionally, one might run a loop until scores are acceptable.
        return new_questions, feedback_text, scores
    else:
        return current_questions, "All evaluation metrics are above thresholds.", scores


# ---------------------------
# Main Code Execution
# ---------------------------
def main():
    pdf_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\GraphAlgorithms.pdf"

    # Initialize Session Manager and Feedback Handler
    session_manager = SessionManager(qdrant_client)
    feedback_handler = FeedbackHandler()

    session_manager.log_interaction("session_start", {"session_id": session_id}) # Log session start

    # ----------- Text Extraction -----------
    print("Starting text extraction (text-only)...\n")
    session_manager.log_interaction("text_extraction_start", {"file_path": pdf_path})
    text_content = extract_text(pdf_path)
    cleaned_text = clean_text(text_content)
    print("Text extraction completed.\n")
    session_manager.log_interaction("text_extraction_complete", {"extracted_text_length": len(cleaned_text)})

    # ----------- Table Extraction -----------
    print("Starting table extraction...\n")
    session_manager.log_interaction("table_extraction_start", {"file_path": pdf_path})
    tables = extract_tables(pdf_path)
    print("Table extraction completed.\n")
    session_manager.log_interaction("table_extraction_complete", {"extracted_table_count": len(tables)})

    # ----------- Diagram/Graph/Illustration Extraction & Insight Analysis -----------
    print("Starting diagram extraction and insight analysis...\n")
    session_manager.log_interaction("diagram_insight_extraction_start", {"file_path": pdf_path})
    diagram_insights = process_file_content_image(pdf_path)
    cleaned_diagram_insights = clean_text(diagram_insights)
    session_manager.log_interaction("diagram_insight_extraction_complete", {"extracted_insights_length": len(cleaned_diagram_insights)})

    # Combine the OCR text with diagram insights.
    combined_text = cleaned_text + " " + cleaned_diagram_insights
    print(f"Combined text length: {len(combined_text)} characters.\n")


    # --- Embedding & Storing in Qdrant (Text Collection Only) -----------
    print("Splitting combined text into chunks...\n")
    session_manager.log_interaction("text_chunking_start", {"text_length": len(combined_text)})
    text_chunks = split_text(combined_text, chunk_size=1024, tolerance=20)  # Use split_text
    print("Embedding text chunks...\n")
    session_manager.log_interaction("embedding_start", {"num_chunks": len(text_chunks)})
    text_embeddings = embed_text_chunks(text_chunks, model)
    print("Inserting text chunks into Qdrant...\n")
    session_manager.log_interaction("qdrant_insertion_start", {"collection_name": "qgen", "num_embeddings": len(text_embeddings)})
    insert_into_qdrant("qgen", text_embeddings, text_chunks, document_metadata)
    print("Text and diagram insights processing and storage completed.\n")
    session_manager.log_interaction("qdrant_insertion_complete", {"collection_name": "qgen", "status": "success"})

    # --- Define parameters (can be from user input or config) ---
    course_name = "Data Structures and Algorithms"
    num_questions = "15"
    academic_level = "undergraduate"
    taxonomy_level = "Analyze"
    topics_list = "BFS, DFS"  # Comma-separated
    major = "Computer Science"
    topics = topics_list  # For the hypothetical text generation
    llm="deepseek"
     # --- Define Thresholds ---
    thresholds = {
    "BLEU-4": 0.1,
    "QSTS": 0.3,
    }


    # --- Hypothetical Sub-topic Generation ---
    hypothetical_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\hypothetical_prompt.txt"
    updated_hypothetical_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\updated_hypothetical_prompt.txt"

    hypothetical_text = find_topics_and_generate_hypothetical_text(hypothetical_prompt_path, updated_hypothetical_path, academic_level, major, course_name, taxonomy_level, topics,llm)
    print("Generated Hypothetical Text (Sub-topics):\n")
    print(hypothetical_text)

    # --- Search Qdrant using Hypothetical Text ---
    print(f"Performing a search in Qdrant for Hypothetical Text...\n")
    query_embedding = model.encode(hypothetical_text)
     # Pass session_id to filter results
    search_results = search_results_from_qdrant(qdrant_client, "qgen", query_embedding, limit=30, session_id_filter=session_id)
    # print("Search Results:")
    retrieved_context = ""
    for result in search_results:
        text_snippet = result.payload.get("text", "")
        # print(f"ID: {result.id}, Score: {result.score:.4f}, Snippet: {text_snippet[:200]}")  # Show first 200 chars
        retrieved_context += text_snippet + " "  # Accumulate context

    # --- Question Generation ---

    user_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\user_prompt.txt"  # Original user prompt
    updated_user_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\updated_user_prompt.txt" # Updated Original user prompt
    final_user_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\final_user_prompt.txt" # Final user prompt
    updated_final_user_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\updated_final_user_prompt.txt" #Updated final user prompt
    evaluation_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\evaluation_prompt.txt"
    updated_evaluation_prompt_path = r"D:\Capstone\CodeFolder\Bloom-s-Question-Generation-using-LLM-s\content\updated_evaluation_prompt.txt"

    #First fill the original user prompt
    placeholders_user = {
        "course_name": course_name,
        "num_questions": num_questions,
        "academic_level": academic_level,
        "topics_list": topics_list,  # This is correct
        "taxonomy": taxonomy_level,
        "major": major
    }
    fill_placeholders(user_prompt_path, updated_user_prompt_path, placeholders_user)

    with open(updated_user_prompt_path, "r", encoding="utf8") as file:
        updated_user_prompt = file.read() #Storing filled user prompt to use it later for final prompt

    # Load the final user prompt content *before* passing to generate_questions, as it's used as base prompt for feedback iterations.
    with open(final_user_prompt_path, "r", encoding="utf8") as file:
        final_user_prompt_content = file.read()


    generated_questions, evaluation_feedback, scores  = generate_questions(
        final_user_prompt_path, updated_final_user_prompt_path, retrieved_context,
        course_name, num_questions, academic_level, taxonomy_level,
        topics_list, major, evaluation_prompt_path, updated_evaluation_prompt_path, thresholds, feedback_handler, final_user_prompt_content,llm # Pass feedback handler and final user prompt content
    )

    print("\nFinal Generated Questions:")
    print(generated_questions)
    print("\nEvaluation Feedback:")
    print(evaluation_feedback)
    print("\nEvaluation scores:")
    print(scores)
    session_manager.log_interaction("question_generation_complete", {"evaluation_scores": scores, "feedback": evaluation_feedback}) # Log question generation completion

    session_manager.log_interaction("session_end", {"session_id": session_id}) # Log session end


if __name__ == "__main__":
    main()